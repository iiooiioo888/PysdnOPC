"""附件內容提取與多模態路由輔助模組。

職責說明：
    負責從使用者上傳的附件中提取文字內容，支援純文字檔案、
    Microsoft Office 格式（.docx、.xlsx、.pptx）。提取的文字用於：
    1. 作為 LLM 上下文的一部分（多模態路由）
    2. 產生附件預覽摘要
    3. 支援附件內容搜尋

關聯關係：
    - 被 opc/core/attachment_store.py 在儲存附件時調用以提取預覽文字
    - 被 opc/layer1_perception/context_assembler.py 在組裝上下文時調用
    - 被 opc/layer4_tools/ 的工具在處理附件時調用

使用範例：
    from opc.core.attachment_content import extract_attachment_text, can_extract_text
    if can_extract_text("report.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
        text = extract_attachment_text("report.docx", mime, raw_bytes)
"""

from __future__ import annotations  # 啟用延遲型別註解評估

from io import BytesIO  # 標準庫：將 bytes 包裝為類檔案物件，供 Office 解析庫讀取
from pathlib import Path  # 標準庫：解析副檔名，跨平台路徑處理
from typing import Iterable  # 標準庫：型別註解，表示可迭代物件

# 可直接以 UTF-8 解碼的純文字副檔名集合。
# 包含：程式碼檔（.py, .js, .ts 等）、設定檔（.yaml, .toml 等）、
# 文件檔（.md, .txt, .csv）、其他（.sql, .log, .html, .css）。
TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".yaml", ".yml",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css",
    ".xml", ".toml", ".ini", ".cfg", ".log", ".sql",
}

# 需要專用解析庫的 Microsoft Office 副檔名集合。
# .docx → python-docx、.xlsx → openpyxl、.pptx → python-pptx
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


class _PreviewAccumulator:
    """文字預覽累積器，用於控制提取文字的總長度。

    職責說明：
        在逐段提取文字時追蹤已使用的字元數，當達到上限時停止累積。
        避免大型文件產生過長的預覽文字（保護 LLM 上下文窗口）。

    使用範例：
        acc = _PreviewAccumulator(max_chars=4000)
        for paragraph in doc.paragraphs:
            if acc.add(paragraph.text):  # 返回 True 表示已達上限
                break
        preview = acc.render()
    """

    def __init__(self, max_chars: int) -> None:
        """初始化累積器。

        參數：
            max_chars (int)：最大累積字元數。負值會被修正為 0。
        """
        # 最大字元數上限（負值保護）
        self.max_chars = max(0, max_chars)
        # 已累積的文字片段列表
        self.parts: list[str] = []
        # 已使用的字元數（包含片段間的分隔符）
        self.used = 0

    def add(self, text: str) -> bool:
        """新增一段文字到累積器。

        功能：
            將文字正規化（去除首尾空白）後加入累積器。
            若剩餘空間不足則截斷文字。

        參數：
            text (str)：要加入的文字片段。

        返回值：
            bool — True 表示已達到或超過字元上限（呼叫者應停止新增）。
        """
        # 正規化：去除首尾空白，空字串不計入
        normalized = str(text or "").strip()
        if not normalized:
            return False
        # 已達上限，通知呼叫者停止
        if self.used >= self.max_chars:
            return True

        # 計算剩餘可用空間
        available = self.max_chars - self.used
        # 超出空間時截斷
        if len(normalized) > available:
            normalized = normalized[:available].rstrip()
        if not normalized:
            return True

        # 加入片段並更新已使用字元數（+1 為換行分隔符）
        self.parts.append(normalized)
        self.used += len(normalized) + 1
        return self.used >= self.max_chars

    def render(self) -> str:
        """將所有累積的片段合併為單一文字。

        返回值：
            str — 以換行符連接的所有文字片段。
        """
        return "\n".join(self.parts).strip()


def attachment_suffix(filename: str) -> str:
    """取得檔案的小寫副檔名。

    功能：
        從檔名中提取副檔名並轉為小寫，用於後續的類型判斷。

    參數：
        filename (str)：檔案名稱（可含路徑），例如 "Report.DOCX"。

    返回值：
        str — 小寫副檔名（含點號），例如 ".docx"。無副檔名時返回空字串。

    被誰引用：
        - is_text_like_attachment()：判斷是否為文字檔
        - can_extract_text()：判斷是否可提取文字
        - extract_attachment_text()：路由到對應的提取函數
    """
    return Path(filename).suffix.lower()


def is_text_like_attachment(filename: str, mime_type: str) -> bool:
    """判斷附件是否為純文字類型（可直接 UTF-8 解碼）。

    功能：
        根據 MIME 類型或副檔名判斷附件是否為純文字格式。
        MIME 以 "text/" 開頭或副檔名在 TEXT_EXTENSIONS 中均視為文字。

    參數：
        filename (str)：檔案名稱，用於提取副檔名。
        mime_type (str)：MIME 類型字串，例如 "text/plain"、"application/json"。

    返回值：
        bool — True 表示為純文字類型，可直接解碼。

    被誰引用：
        - can_extract_text()：作為判斷條件之一
        - extract_attachment_text()：決定使用直接解碼路徑
    """
    if mime_type.startswith("text/"):
        return True
    return attachment_suffix(filename) in TEXT_EXTENSIONS


def can_extract_text(filename: str, mime_type: str) -> bool:
    """判斷附件是否可以提取文字內容。

    功能：
        綜合判斷純文字檔和 Office 文件，決定是否能從中提取有意義的文字。
        用於在處理附件前快速判斷是否需要進行內容提取。

    參數：
        filename (str)：檔案名稱。
        mime_type (str)：MIME 類型字串。

    返回值：
        bool — True 表示可以提取文字（純文字或 Office 格式）。

    被誰引用：
        - opc/core/attachment_store.py：儲存前判斷是否提取預覽
        - opc/layer1_perception/context_assembler.py：組裝上下文時過濾
    """
    return is_text_like_attachment(filename, mime_type) or attachment_suffix(filename) in OFFICE_EXTENSIONS


def extract_attachment_text(
    filename: str,
    mime_type: str,
    raw: bytes,
    *,
    max_chars: int = 4000,
) -> str:
    """從附件中提取文字內容（主入口函數）。

    功能：
        根據檔案類型路由到對應的提取邏輯：
        - 純文字檔：直接 UTF-8 解碼
        - .docx：使用 python-docx 提取段落和表格
        - .xlsx：使用 openpyxl 提取工作表儲存格
        - .pptx：使用 python-pptx 提取投影片文字
        所有提取結果均截斷至 max_chars 字元。

    參數：
        filename (str)：檔案名稱（含副檔名），用於判斷類型。
        mime_type (str)：MIME 類型，輔助判斷純文字檔。
        raw (bytes)：檔案的原始位元組內容。
        max_chars (int)：提取文字的最大字元數。預設 4000。
            取值範圍：正整數，過大可能影響 LLM 上下文效能。

    返回值：
        str — 提取的文字內容（可能已截斷）。無法提取時返回空字串。

    被誰引用：
        - opc/core/attachment_store.py：儲存附件時產生預覽
        - opc/layer1_perception/context_assembler.py：將附件內容加入 LLM 上下文
    """
    suffix = attachment_suffix(filename)
    # 純文字檔：直接 UTF-8 解碼（errors="replace" 處理非法位元組）
    if is_text_like_attachment(filename, mime_type):
        return _clip_text(raw.decode("utf-8", errors="replace").strip(), max_chars)
    # Word 文件
    if suffix == ".docx":
        return _extract_docx_text(raw, max_chars=max_chars)
    # Excel 試算表
    if suffix == ".xlsx":
        return _extract_xlsx_text(raw, max_chars=max_chars)
    # PowerPoint 簡報
    if suffix == ".pptx":
        return _extract_pptx_text(raw, max_chars=max_chars)
    # 不支援的格式返回空字串
    return ""


def _clip_text(text: str, max_chars: int) -> str:
    """將文字截斷至指定長度，超出部分以省略標記取代。

    參數：
        text (str)：原始文字。
        max_chars (int)：最大允許字元數。

    返回值：
        str — 截斷後的文字。若未超出則原樣返回。
    """
    text = str(text or "").strip()
    if len(text) <= max_chars:
        return text
    return f"{text[:max_chars].rstrip()}\n...[truncated]"


def _extract_docx_text(raw: bytes, *, max_chars: int) -> str:
    """從 .docx 檔案中提取文字（段落 + 表格）。

    功能：
        使用 python-docx 庫解析 Word 文件，依序提取所有段落文字
        和表格內容，累積至 max_chars 上限。

    參數：
        raw (bytes)：.docx 檔案的原始位元組。
        max_chars (int)：最大提取字元數。

    返回值：
        str — 提取的文字內容。

    依賴：
        python-docx（延遲導入，僅在處理 .docx 時載入）
    """
    from docx import Document  # 延遲導入：python-docx 庫，解析 .docx 格式

    acc = _PreviewAccumulator(max_chars)
    doc = Document(BytesIO(raw))

    # 第一優先：提取所有段落文字
    for para in doc.paragraphs:
        if acc.add(para.text):
            return _clip_text(acc.render(), max_chars)

    # 第二優先：提取表格內容（以 " | " 分隔儲存格）
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if acc.add(line):
                return _clip_text(acc.render(), max_chars)

    return _clip_text(acc.render(), max_chars)


def _extract_xlsx_text(raw: bytes, *, max_chars: int) -> str:
    """從 .xlsx 檔案中提取文字（前 5 個工作表，每表最多 80 行）。

    功能：
        使用 openpyxl 庫以唯讀模式解析 Excel 檔案，提取前 5 個
        工作表的內容。每行最多取前 16 個儲存格，以定位符分隔。

    參數：
        raw (bytes)：.xlsx 檔案的原始位元組。
        max_chars (int)：最大提取字元數。

    返回值：
        str — 提取的文字內容，包含工作表標題標記。

    依賴：
        openpyxl（延遲導入，僅在處理 .xlsx 時載入）
    """
    from openpyxl import load_workbook  # 延遲導入：openpyxl 庫，解析 .xlsx 格式

    acc = _PreviewAccumulator(max_chars)
    # read_only=True：唯讀模式，減少記憶體佔用；data_only=True：讀取計算結果而非公式
    workbook = load_workbook(BytesIO(raw), read_only=True, data_only=True)
    try:
        # 最多處理前 5 個工作表
        for sheet in workbook.worksheets[:5]:
            if acc.add(f"# Sheet: {sheet.title}"):
                break
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                # 每行最多取前 16 個儲存格
                values = [_normalize_excel_cell(value) for value in row[:16]]
                # 跳過全空行
                if not any(values):
                    continue
                row_count += 1
                if acc.add("\t".join(values)):
                    return _clip_text(acc.render(), max_chars)
                # 每個工作表最多 80 行
                if row_count >= 80:
                    break
    finally:
        workbook.close()  # 確保釋放檔案資源

    return _clip_text(acc.render(), max_chars)


def _normalize_excel_cell(value: object) -> str:
    """將 Excel 儲存格值正規化為字串。

    功能：
        處理 None、浮點數（去除多餘的尾零）、一般值的字串轉換。

    參數：
        value (object)：儲存格的原始值（可能為 None、float、str 等）。

    返回值：
        str — 正規化後的字串。None → 空字串；浮點數去除尾零。
    """
    if value is None:
        return ""
    if isinstance(value, float):
        # 格式化為 6 位小數後去除尾零（3.140000 → 3.14）
        text = f"{value:.6f}".rstrip("0").rstrip(".")
        return text or "0"
    return str(value).strip()


def _extract_pptx_text(raw: bytes, *, max_chars: int) -> str:
    """從 .pptx 檔案中提取文字（前 20 張投影片）。

    功能：
        使用 python-pptx 庫解析 PowerPoint 檔案，提取前 20 張
        投影片中的所有文字框和表格內容。

    參數：
        raw (bytes)：.pptx 檔案的原始位元組。
        max_chars (int)：最大提取字元數。

    返回值：
        str — 提取的文字內容，包含投影片編號標記。

    依賴：
        python-pptx（延遲導入，僅在處理 .pptx 時載入）
    """
    from pptx import Presentation  # 延遲導入：python-pptx 庫，解析 .pptx 格式

    acc = _PreviewAccumulator(max_chars)
    presentation = Presentation(BytesIO(raw))

    # 最多處理前 20 張投影片
    for index, slide in enumerate(list(presentation.slides)[:20], start=1):
        if acc.add(f"# Slide {index}"):
            break
        for text in _iter_slide_text(slide.shapes):
            if acc.add(text):
                return _clip_text(acc.render(), max_chars)

    return _clip_text(acc.render(), max_chars)


def _iter_slide_text(shapes: Iterable[object]) -> Iterable[str]:
    """遞迴迭代投影片中的所有文字內容（生成器）。

    功能：
        遍歷投影片圖形集合，提取：
        1. 文字框的直接文字
        2. 表格中的儲存格文字（以 " | " 分隔）
        3. 群組圖形中的子圖形文字（遞迴處理）

    參數：
        shapes (Iterable[object])：投影片的圖形集合。

    返回值：
        Iterable[str] — 依序產生的文字字串（生成器）。

    被誰引用：
        - _extract_pptx_text()：提取投影片文字
    """
    for shape in shapes:
        # 提取文字框內容
        text = getattr(shape, "text", "")
        if isinstance(text, str) and text.strip():
            yield text

        # 提取表格內容
        table = getattr(shape, "table", None)
        if table is not None:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    yield line

        # 遞迴處理群組圖形中的子圖形
        subshapes = getattr(shape, "shapes", None)
        if subshapes is not None:
            yield from _iter_slide_text(subshapes)
